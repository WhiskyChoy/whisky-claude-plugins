#!/usr/bin/env python3
"""
diff_supplement.py — Extract and optionally diff supplemental materials against
existing HTML slides.

Two modes:
  --extract-only  Extract structured content from all inputs as JSON for
                  semantic analysis by the LLM (recommended).
  (default)       Syntactic fuzzy-match via SequenceMatcher. Use only as a
                  rough hint — the LLM should make final merge decisions.

Usage:
    # Recommended: extract content for LLM semantic comparison
    python diff_supplement.py slides.html supplement1.pptx --extract-only -o content.json

    # Legacy: syntactic diff (use as auxiliary hint only)
    python diff_supplement.py slides.html supplement1.pptx -o merge_plan.json

Dependencies (auto-installed if missing):
    beautifulsoup4, lxml, pdfplumber, python-pptx
"""

import argparse
import difflib
import json
import os
import subprocess
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional


# --- Dependency management ---------------------------------------------------

def ensure_import(package_name, pip_name=None):
    """Import a package, auto-installing via pip if missing."""
    try:
        return __import__(package_name)
    except ImportError:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", pip_name or package_name],
            stdout=subprocess.DEVNULL,
        )
        return __import__(package_name)


# --- Data structures ---------------------------------------------------------

@dataclass
class SlideContent:
    """Represents the content of a single slide or section."""
    index: int
    heading: str
    body_text: str
    figures: list = field(default_factory=list)
    slide_type: str = "content"  # title/content/figure/section/end
    source_file: str = ""

    def text_for_matching(self):
        """Return combined text for fuzzy matching."""
        return f"{self.heading}\n{self.body_text}".strip()


@dataclass
class DiffEntry:
    """A single entry in the merge plan."""
    action: str  # "new", "update", "unchanged"
    supplement_index: int
    supplement_heading: str
    supplement_body_preview: str
    existing_index: Optional[int] = None
    similarity: float = 0.0
    source_file: str = ""


# --- Content extraction ------------------------------------------------------

def extract_from_html(html_path):
    """Extract slide content from an HTML presentation file.

    Returns a list of SlideContent.
    """
    ensure_import("bs4", "beautifulsoup4")
    ensure_import("lxml")
    from bs4 import BeautifulSoup

    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "lxml")

    slides = []
    slide_elements = (
        soup.select("section.slide")
        or soup.select("div.slide")
        or soup.select("section")
        or soup.select("[data-slide]")
    )

    for idx, el in enumerate(slide_elements):
        heading_el = el.find(["h1", "h2", "h3"])
        heading = heading_el.get_text(strip=True) if heading_el else ""

        # Collect body text (excluding the heading)
        body_parts = []
        for child in el.find_all(["p", "li", "span", "div"], recursive=True):
            if child.find_parent(["h1", "h2", "h3"]):
                continue
            text = child.get_text(strip=True)
            if text and text != heading:
                body_parts.append(text)

        body_text = "\n".join(dict.fromkeys(body_parts))  # deduplicate order-preserving

        # Detect figures
        figures = [
            img.get("src", "") for img in el.find_all("img")
        ]

        # Classify slide type
        classes = " ".join(el.get("class", []))
        if idx == 0 or "title" in classes:
            stype = "title"
        elif idx == len(slide_elements) - 1 or any(
            kw in classes for kw in ("end", "thank", "final")
        ):
            stype = "end"
        elif figures and not body_text.strip():
            stype = "figure"
        else:
            stype = "content"

        slides.append(SlideContent(
            index=idx,
            heading=heading,
            body_text=body_text[:2000],
            figures=figures,
            slide_type=stype,
            source_file=str(html_path),
        ))

    return slides


def extract_from_pptx(pptx_path):
    """Extract slide content from a PPTX file.

    Returns a list of SlideContent.
    """
    ensure_import("pptx", "python-pptx")
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides = []

    for idx, slide in enumerate(prs.slides):
        texts = []
        heading = ""
        figures = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # First significant text block as heading
                        if not heading and len(text) < 200:
                            heading = text
                        else:
                            texts.append(text)

            if hasattr(shape, "image"):
                try:
                    figures.append(shape.name)
                except Exception:
                    pass

        body_text = "\n".join(texts)

        # Classify
        if idx == 0:
            stype = "title"
        elif idx == len(prs.slides) - 1:
            stype = "end"
        elif figures and not body_text.strip():
            stype = "figure"
        else:
            stype = "content"

        slides.append(SlideContent(
            index=idx,
            heading=heading,
            body_text=body_text[:2000],
            figures=figures,
            slide_type=stype,
            source_file=str(pptx_path),
        ))

    return slides


def extract_from_pdf(pdf_path):
    """Extract page content from a PDF file.

    Returns a list of SlideContent (one per page).
    """
    ensure_import("pdfplumber")
    import pdfplumber                       # type: ignore[import]  # pdfplumber does not have type hints

    slides = []
    with pdfplumber.open(pdf_path) as pdf:
        for idx, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            lines = text.strip().split("\n")

            # First line as heading, rest as body
            heading = lines[0] if lines else ""
            body_text = "\n".join(lines[1:]) if len(lines) > 1 else ""

            slides.append(SlideContent(
                index=idx,
                heading=heading[:200],
                body_text=body_text[:2000],
                slide_type="content",
                source_file=str(pdf_path),
            ))

    return slides


def extract_content(file_path):
    """Dispatch to the correct extractor based on file extension."""
    ext = Path(file_path).suffix.lower()

    extractors = {
        ".html": extract_from_html,
        ".htm": extract_from_html,
        ".pptx": extract_from_pptx,
        ".pdf": extract_from_pdf,
    }

    extractor = extractors.get(ext)
    if not extractor:
        print(f"Warning: Unsupported file type '{ext}' for {file_path}",
              file=sys.stderr)
        return []

    print(f"Extracting content from: {file_path} ({ext})")
    return extractor(file_path)


# --- Diff algorithm ----------------------------------------------------------

MATCH_THRESHOLD = 0.6


def diff_slides(existing, supplements, threshold=MATCH_THRESHOLD):
    """Compare supplemental slides against existing slides using fuzzy matching.

    Args:
        existing: list of SlideContent from the existing slides.html
        supplements: list of SlideContent from supplemental files
        threshold: similarity threshold for matching (default: 0.6)

    Returns:
        A list of DiffEntry describing the merge plan.
    """
    results = []
    matched_existing = set()

    for supp in supplements:
        supp_text = supp.text_for_matching()
        if not supp_text.strip():
            continue

        best_score = 0.0
        best_idx = None

        for ex in existing:
            if ex.index in matched_existing:
                continue
            ex_text = ex.text_for_matching()
            if not ex_text.strip():
                continue

            score = difflib.SequenceMatcher(
                None, supp_text.lower(), ex_text.lower()
            ).ratio()

            if score > best_score:
                best_score = score
                best_idx = ex.index

        body_preview = supp.body_text[:150].replace("\n", " ")

        if best_score >= 0.9:
            # Unchanged — very high match
            matched_existing.add(best_idx)
            results.append(DiffEntry(
                action="unchanged",
                supplement_index=supp.index,
                supplement_heading=supp.heading,
                supplement_body_preview=body_preview,
                existing_index=best_idx,
                similarity=round(best_score, 3),
                source_file=supp.source_file,
            ))
        elif best_score >= threshold:
            # Updated — partial match
            matched_existing.add(best_idx)
            results.append(DiffEntry(
                action="update",
                supplement_index=supp.index,
                supplement_heading=supp.heading,
                supplement_body_preview=body_preview,
                existing_index=best_idx,
                similarity=round(best_score, 3),
                source_file=supp.source_file,
            ))
        else:
            # New — no good match found
            results.append(DiffEntry(
                action="new",
                supplement_index=supp.index,
                supplement_heading=supp.heading,
                supplement_body_preview=body_preview,
                similarity=round(best_score, 3) if best_idx is not None else 0.0,
                source_file=supp.source_file,
            ))

    return results


# --- Main --------------------------------------------------------------------

def run_extract_only(existing_path, supplement_paths, output_path):
    """Extract structured content from all inputs without comparison.

    Outputs a JSON file with full text for each slide/section, designed
    for the LLM to read and perform semantic comparison.
    """
    print(f"\n--- Extracting (extract-only mode) ---")

    existing = extract_content(existing_path)
    print(f"  Existing slides: {len(existing)} extracted")

    all_supplements = []
    for s in supplement_paths:
        content = extract_content(s)
        all_supplements.extend(content)
        print(f"  Supplement {s}: {len(content)} sections extracted")

    output = {
        "mode": "extract-only",
        "note": (
            "This file contains extracted text for semantic comparison by the "
            "LLM. Read both sections and judge similarity by meaning, not "
            "wording. Paraphrased content covering the same topic should be "
            "treated as matching."
        ),
        "existing_slides": {
            "source": str(existing_path),
            "count": len(existing),
            "slides": [
                {
                    "index": s.index,
                    "heading": s.heading,
                    "body_text": s.body_text,
                    "slide_type": s.slide_type,
                    "figures": s.figures,
                }
                for s in existing
            ],
        },
        "supplements": {
            "sources": [str(s) for s in supplement_paths],
            "count": len(all_supplements),
            "sections": [
                {
                    "index": s.index,
                    "heading": s.heading,
                    "body_text": s.body_text,
                    "slide_type": s.slide_type,
                    "figures": s.figures,
                    "source_file": s.source_file,
                }
                for s in all_supplements
            ],
        },
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=2, ensure_ascii=False)

    print(f"\nExtracted content saved to: {output_path}")
    print("  → Read this file and perform semantic comparison yourself.")


def main():
    parser = argparse.ArgumentParser(
        description=(
            "Extract and optionally diff supplemental materials against "
            "existing HTML slides. Use --extract-only (recommended) for "
            "LLM-based semantic comparison."
        ),
    )
    parser.add_argument(
        "existing_slides",
        help="Path to the existing slides.html file",
    )
    parser.add_argument(
        "supplements",
        nargs="+",
        help="Path(s) to supplemental files (.pptx, .pdf, .html)",
    )
    parser.add_argument(
        "--output", "-o",
        default="merge_plan.json",
        help="Output path for the JSON file (default: merge_plan.json)",
    )
    parser.add_argument(
        "--extract-only", "-e",
        action="store_true",
        help=(
            "Extract content only — no comparison. Outputs structured text "
            "for the LLM to perform semantic analysis (recommended)."
        ),
    )
    parser.add_argument(
        "--threshold", "-t",
        type=float,
        default=MATCH_THRESHOLD,
        help=f"Similarity threshold for syntactic matching (default: {MATCH_THRESHOLD})",
    )
    args = parser.parse_args()

    # Validate inputs
    if not os.path.isfile(args.existing_slides):
        print(f"Error: existing slides not found: {args.existing_slides}",
              file=sys.stderr)
        sys.exit(1)

    for s in args.supplements:
        if not os.path.isfile(s):
            print(f"Error: supplement not found: {s}", file=sys.stderr)
            sys.exit(1)

    # Extract-only mode: output content for LLM semantic comparison
    if args.extract_only:
        run_extract_only(args.existing_slides, args.supplements, args.output)
        return

    # --- Legacy syntactic diff mode ---
    print(
        "\nNote: Syntactic diff uses string similarity (SequenceMatcher). "
        "It may misjudge semantically similar content as different. "
        "Use --extract-only for LLM-based semantic comparison instead.\n"
    )

    # Extract existing slides
    print(f"--- Existing slides ---")
    existing = extract_content(args.existing_slides)
    print(f"  {len(existing)} slides extracted\n")

    # Extract all supplements
    all_supplements = []
    for s in args.supplements:
        print(f"--- Supplement: {s} ---")
        content = extract_content(s)
        all_supplements.extend(content)
        print(f"  {len(content)} sections extracted\n")

    # Diff
    print("--- Running syntactic diff (auxiliary hint only) ---")
    results = diff_slides(existing, all_supplements, threshold=args.threshold)

    # Summarize
    new_count = sum(1 for r in results if r.action == "new")
    update_count = sum(1 for r in results if r.action == "update")
    unchanged_count = sum(1 for r in results if r.action == "unchanged")

    print(f"\nResults (syntactic — verify semantically before acting):")
    print(f"  New sections:       {new_count}")
    print(f"  Updated sections:   {update_count}")
    print(f"  Unchanged sections: {unchanged_count}")

    # Build output
    merge_plan = {
        "existing_slides": str(args.existing_slides),
        "supplements": [str(s) for s in args.supplements],
        "existing_slide_count": len(existing),
        "supplement_section_count": len(all_supplements),
        "threshold": args.threshold,
        "warning": (
            "This diff is based on string similarity (SequenceMatcher), not "
            "semantic meaning. Paraphrased or restructured content may be "
            "incorrectly classified as 'new'. Always verify by reading the "
            "actual content before making merge decisions."
        ),
        "summary": {
            "new": new_count,
            "updated": update_count,
            "unchanged": unchanged_count,
        },
        "entries": [asdict(r) for r in results],
    }

    # Write output
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(merge_plan, f, indent=2, ensure_ascii=False)

    print(f"\nMerge plan saved to: {args.output}")


if __name__ == "__main__":
    main()
