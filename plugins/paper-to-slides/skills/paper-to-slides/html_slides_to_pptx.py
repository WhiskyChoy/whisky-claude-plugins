#!/usr/bin/env python3
"""
html_slides_to_pptx.py — Pixel-perfect HTML slides → PPTX via Playwright screenshots.

Renders each slide in a headless browser at the target resolution, then embeds
each screenshot as a full-slide image in a PPTX file. Guarantees 100% visual
fidelity: fonts, math, backgrounds, logos all captured exactly as rendered.

Usage:
    python html_slides_to_pptx.py slides.html --output slides.pptx
    python html_slides_to_pptx.py slides.html --output slides.pptx --ratio 4:3
    python html_slides_to_pptx.py slides.html --output slides.pptx --ratio 16:9 --dpi 192
"""

from __future__ import annotations

import argparse
import subprocess
import sys
import time
from pathlib import Path


def ensure_deps():
    """Auto-install missing dependencies."""
    missing = []
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError:
        missing.append("playwright")
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        missing.append("python-pptx")

    if missing:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", *missing],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    # Ensure Chromium is installed
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            p.chromium.launch(headless=True).close()
    except Exception:
        subprocess.check_call(
            [sys.executable, "-m", "playwright", "install", "chromium"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )


# Aspect ratio presets: (width_inches, height_inches, viewport_w, viewport_h)
RATIO_PRESETS = {
    "16:9": (13.333, 7.5, 1920, 1080),
    "4:3":  (10.0,   7.5, 1440, 1080),
}


def render_slides(
    html_path: Path,
    screenshot_dir: Path,
    viewport_w: int,
    viewport_h: int,
    dpi_scale: float,
) -> list[Path]:
    """Render each slide in the HTML to a PNG screenshot."""
    from playwright.sync_api import sync_playwright

    screenshot_dir.mkdir(parents=True, exist_ok=True)
    file_url = html_path.resolve().as_uri()
    screenshots = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            viewport={"width": viewport_w, "height": viewport_h},
            device_scale_factor=dpi_scale,
        )

        # Load the HTML file
        page.goto(file_url, wait_until="domcontentloaded", timeout=60000)

        # Wait for KaTeX to finish rendering
        page.wait_for_timeout(1500)

        # Count slides
        total = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Found {total} slides. Rendering at {viewport_w}x{viewport_h} "
              f"(scale {dpi_scale}x)...")

        for i in range(total):
            # Activate the i-th slide, deactivate others
            page.evaluate(f"""(() => {{
                const slides = document.querySelectorAll('.slide');
                slides.forEach((s, idx) => {{
                    s.classList.toggle('active', idx === {i});
                }});
                // Force reveal-items to be visible immediately
                const active = slides[{i}];
                active.querySelectorAll('.reveal-item').forEach(el => {{
                    el.style.opacity = '1';
                    el.style.transform = 'none';
                    el.style.transition = 'none';
                }});
            }})()""")

            # Brief pause for rendering
            page.wait_for_timeout(300)

            png_path = screenshot_dir / f"slide_{i + 1:03d}.png"
            page.screenshot(path=str(png_path))
            screenshots.append(png_path)
            print(f"  Slide {i + 1}/{total}")

        browser.close()

    return screenshots


def assemble_pptx(
    screenshots: list[Path],
    output_path: Path,
    width_inches: float,
    height_inches: float,
):
    """Assemble screenshots into a PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)

    blank_layout = prs.slide_layouts[6]  # Blank

    for png_path in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output_path))
    print(f"\nSaved PPTX: {output_path} ({len(screenshots)} slides)")


def main():
    parser = argparse.ArgumentParser(
        description="Pixel-perfect HTML slides → PPTX via browser screenshots."
    )
    parser.add_argument("html_path", help="Path to the HTML slides file")
    parser.add_argument(
        "--output", "-o", required=True,
        help="Output PPTX file path"
    )
    parser.add_argument(
        "--ratio", choices=["16:9", "4:3"], default="16:9",
        help="Slide aspect ratio (default: 16:9)"
    )
    parser.add_argument(
        "--dpi", type=int, default=2,
        help="DPI scale factor for rendering (default: 2 for crisp output)"
    )
    parser.add_argument(
        "--keep-screenshots", action="store_true",
        help="Keep the intermediate screenshot PNGs"
    )

    args = parser.parse_args()
    html_path = Path(args.html_path)
    output_path = Path(args.output)

    if not html_path.exists():
        print(f"Error: {html_path} not found")
        sys.exit(1)

    width_in, height_in, vw, vh = RATIO_PRESETS[args.ratio]
    screenshot_dir = output_path.resolve().parent / f"_pptx_screenshots_{output_path.stem}"

    ensure_deps()

    t0 = time.time()
    screenshots = render_slides(html_path, screenshot_dir, vw, vh, args.dpi)
    assemble_pptx(screenshots, output_path, width_in, height_in)
    elapsed = time.time() - t0

    # Cleanup screenshots unless asked to keep
    if not args.keep_screenshots:
        for f in screenshots:
            f.unlink(missing_ok=True)
        screenshot_dir.rmdir()
        print("Cleaned up screenshots.")

    print(f"Done in {elapsed:.1f}s")


if __name__ == "__main__":
    main()
