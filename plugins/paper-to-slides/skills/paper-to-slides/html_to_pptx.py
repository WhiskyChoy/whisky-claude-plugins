#!/usr/bin/env python3
"""
html_to_pptx.py - Convert HTML slide presentations to PPTX with proper
math formula rendering, figure fitting, and background support.

Fixes common issues:
  - Math formulas rendered as images (not raw text with "_" for subscripts)
  - Figures auto-scaled to fit within slide boundaries
  - Background images properly applied to slides
  - Logos positioned correctly from template assets

Usage:
    python html_to_pptx.py slides.html --output slides.pptx
    python html_to_pptx.py slides.html --output slides.pptx --assets ./assets
    python html_to_pptx.py slides.html --output slides.pptx --width 13.333 --height 7.5

Dependencies (auto-installed if missing):
    python-pptx, beautifulsoup4, lxml, matplotlib, Pillow
"""

import argparse
import io
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path


def ensure_package(package_name, import_name=None):
    """Auto-install a package if not available."""
    import_name = import_name or package_name
    try:
        __import__(import_name)
    except ImportError:
        subprocess.check_call(
            [sys.executable, '-m', 'pip', 'install', package_name],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )


# Auto-install dependencies
ensure_package('python-pptx', 'pptx')
ensure_package('beautifulsoup4', 'bs4')
ensure_package('lxml')
ensure_package('matplotlib')
ensure_package('Pillow', 'PIL')

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
from PIL import Image


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Default 16:9 widescreen slide dimensions (inches)
DEFAULT_WIDTH = 13.333
DEFAULT_HEIGHT = 7.5

# Margins (inches)
MARGIN_LEFT = 0.7
MARGIN_RIGHT = 0.7
MARGIN_TOP = 0.9
MARGIN_BOTTOM = 0.5

# Math rendering
MATH_DPI = 300
MATH_FONTSIZE = 18
MATH_COLOR = '#333333'

# Regex patterns for LaTeX math in HTML
# Matches KaTeX/MathJax rendered spans and raw LaTeX delimiters
LATEX_INLINE_PATTERN = re.compile(
    r'(?:'
    r'\$([^$]+)\$'            # $...$
    r'|\\\\?\((.+?)\\\\?\)'   # \(...\) or \\(...\\)
    r')',
    re.DOTALL
)
LATEX_DISPLAY_PATTERN = re.compile(
    r'(?:'
    r'\$\$([^$]+)\$\$'         # $$...$$
    r'|\\\\?\[(.+?)\\\\?\]'    # \[...\] or \\[...\\]
    r')',
    re.DOTALL
)


# ---------------------------------------------------------------------------
# Math Rendering
# ---------------------------------------------------------------------------

def render_latex_to_png(latex_str, dpi=MATH_DPI, fontsize=MATH_FONTSIZE,
                        color=MATH_COLOR, display_mode=False):
    """Render a LaTeX math string to a PNG image bytes buffer.

    Returns (png_bytes, width_inches, height_inches) or None on failure.
    """
    latex_str = latex_str.strip()
    if not latex_str:
        return None

    # Ensure the string is wrapped in $ for matplotlib
    if not latex_str.startswith('$'):
        latex_str = f'${latex_str}$'

    try:
        fig = plt.figure(figsize=(0.01, 0.01))
        fig.patch.set_alpha(0)

        text = fig.text(
            0, 0, latex_str,
            fontsize=fontsize if not display_mode else fontsize + 4,
            color=color,
            usetex=False,  # Use mathtext (no LaTeX install required)
        )

        # Render to get bounding box
        fig.canvas.draw()
        bbox = text.get_window_extent(fig.canvas.get_renderer())
        bbox_inches = bbox.transformed(fig.dpi_scale_trans.inverted())

        # Add small padding
        pad = 0.05
        new_w = bbox_inches.width + 2 * pad
        new_h = bbox_inches.height + 2 * pad
        fig.set_size_inches(new_w, new_h)

        # Reposition text
        text.set_position((pad / new_w, pad / new_h))

        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=dpi, transparent=True,
                    bbox_inches='tight', pad_inches=0.02)
        plt.close(fig)

        buf.seek(0)
        # Get actual image dimensions
        with Image.open(buf) as img:
            w_px, h_px = img.size
        buf.seek(0)

        return buf, w_px / dpi, h_px / dpi
    except Exception as e:
        plt.close('all')
        print(f"  Warning: Failed to render math '{latex_str[:50]}...': {e}",
              file=sys.stderr)
        return None


def extract_math_from_katex_span(element):
    """Extract LaTeX source from a KaTeX-rendered span element.

    KaTeX renders into <span class="katex"> with the original LaTeX in
    <annotation encoding="application/x-tex">.
    """
    annotation = element.find('annotation', encoding='application/x-tex')
    if annotation:
        return annotation.get_text()

    # Fallback: try aria-label
    label = element.get('aria-label', '')
    if label:
        return label

    return None


def extract_math_from_mathjax(element):
    """Extract LaTeX from MathJax rendered elements."""
    # MathJax v3: <mjx-container> with a <script type="math/tex">
    script = element.find('script', type=re.compile(r'math/tex'))
    if script:
        return script.get_text()

    # Fallback to alt text on MathJax img
    alt = element.get('alt', '')
    if alt:
        return alt

    return None


# ---------------------------------------------------------------------------
# Image Fitting
# ---------------------------------------------------------------------------

def fit_image_to_area(img_path, max_width_in, max_height_in):
    """Calculate dimensions to fit image within a bounding box, preserving
    aspect ratio.

    Returns (width_inches, height_inches).
    """
    try:
        with Image.open(img_path) as img:
            w_px, h_px = img.size
    except Exception:
        return max_width_in, max_height_in

    if w_px == 0 or h_px == 0:
        return max_width_in, max_height_in

    aspect = w_px / h_px
    target_aspect = max_width_in / max_height_in

    if aspect > target_aspect:
        # Image is wider than area — fit to width
        return max_width_in, max_width_in / aspect
    else:
        # Image is taller than area — fit to height
        return max_height_in * aspect, max_height_in


def add_background_image(slide, img_path, slide_width, slide_height):
    """Add an image as a full-slide background (behind all other shapes)."""
    try:
        pic = slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            slide_width, slide_height
        )
        # Move to back (index 0)
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)  # After background shape
    except Exception as e:
        print(f"  Warning: Could not add background image: {e}",
              file=sys.stderr)


def add_logo(slide, logo_path, position='top-right',
             logo_width_in=0.8, slide_width_in=DEFAULT_WIDTH,
             slide_height_in=DEFAULT_HEIGHT,
             exact_left_in=None, exact_top_in=None):
    """Add a logo image at the specified position.

    If exact_left_in and exact_top_in are provided (from style_report.json),
    use those exact coordinates instead of the named position.
    """
    if exact_left_in is not None and exact_top_in is not None:
        left, top = exact_left_in, exact_top_in
    else:
        positions = {
            'top-right': (slide_width_in - logo_width_in - 0.3, 0.2),
            'top-left': (0.3, 0.2),
            'bottom-right': (slide_width_in - logo_width_in - 0.3,
                             slide_height_in - logo_width_in - 0.2),
            'bottom-left': (0.3, slide_height_in - logo_width_in - 0.2),
        }
        left, top = positions.get(position, positions['top-right'])

    try:
        slide.shapes.add_picture(
            logo_path,
            Inches(left), Inches(top),
            Inches(logo_width_in)
        )
    except Exception as e:
        print(f"  Warning: Could not add logo: {e}", file=sys.stderr)


# ---------------------------------------------------------------------------
# HTML Parsing
# ---------------------------------------------------------------------------

def parse_html_slides(html_path):
    """Parse an HTML presentation file and extract slide data.

    Returns a list of slide dicts:
    [
        {
            'type': 'title' | 'content' | 'end',
            'heading': str,
            'subheading': str | None,
            'body_parts': [
                {'kind': 'text', 'content': str},
                {'kind': 'math_inline', 'latex': str},
                {'kind': 'math_display', 'latex': str},
                {'kind': 'image', 'src': str, 'alt': str},
                {'kind': 'bullet', 'items': [str]},
            ],
            'background_image': str | None,
            'notes': str | None,
        },
        ...
    ]
    """
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    slides = []
    html_dir = Path(html_path).parent

    # Find slide containers — common patterns in HTML presentations
    slide_elements = (
        soup.select('section.slide') or
        soup.select('div.slide') or
        soup.select('section') or
        soup.select('[data-slide]')
    )

    if not slide_elements:
        print("Warning: No slide elements found in HTML. "
              "Looking for <section> or <div class='slide'>.",
              file=sys.stderr)
        return slides

    for idx, el in enumerate(slide_elements):
        slide_data = {
            'type': 'content',
            'heading': '',
            'subheading': None,
            'body_parts': [],
            'background_image': None,
            'notes': None,
        }

        # Detect slide type
        classes = ' '.join(el.get('class', []))
        if idx == 0 or 'title' in classes:
            slide_data['type'] = 'title'
        elif idx == len(slide_elements) - 1 or any(
            kw in classes for kw in ['end', 'thank', 'final']
        ):
            slide_data['type'] = 'end'

        # Extract heading
        heading_el = el.find(['h1', 'h2', 'h3'])
        if heading_el:
            slide_data['heading'] = heading_el.get_text(strip=True)

        # Extract subheading
        if slide_data['type'] == 'title':
            sub_el = el.find(['h2', 'h3', 'p'], class_=re.compile(
                r'sub|author|affil|date'
            ))
            if sub_el and sub_el != heading_el:
                slide_data['subheading'] = sub_el.get_text(strip=True)

        # Extract background image from inline style
        style = el.get('style', '')
        bg_match = re.search(r"background-image:\s*url\(['\"]?([^)'\"]+ )['\"]?\)", style)
        if bg_match:
            bg_src = bg_match.group(1)
            bg_path = html_dir / bg_src
            if bg_path.exists():
                slide_data['background_image'] = str(bg_path)

        # Extract body content
        body_parts = []

        # Process all content elements
        for child in el.descendants:
            if child.name == 'li':
                # Collect bullet items from this list
                parent_list = child.parent
                if parent_list and parent_list.name in ('ul', 'ol'):
                    # Only process the list once (via its first li)
                    if child == parent_list.find('li'):
                        items = [
                            _clean_text_with_math(li)
                            for li in parent_list.find_all('li', recursive=False)
                        ]
                        body_parts.append({'kind': 'bullet', 'items': items})

            elif child.name == 'img':
                src = child.get('src', '')
                img_path = html_dir / src
                if img_path.exists():
                    body_parts.append({
                        'kind': 'image',
                        'src': str(img_path),
                        'alt': child.get('alt', ''),
                    })

            elif child.name == 'span' and 'katex' in ' '.join(
                child.get('class', [])
            ):
                latex = extract_math_from_katex_span(child)
                if latex:
                    is_display = 'katex-display' in ' '.join(
                        child.get('class', [])
                    )
                    body_parts.append({
                        'kind': 'math_display' if is_display else 'math_inline',
                        'latex': latex,
                    })

            elif child.name and child.get('class') and any(
                'MathJax' in c for c in child.get('class', [])
            ):
                latex = extract_math_from_mathjax(child)
                if latex:
                    body_parts.append({
                        'kind': 'math_display',
                        'latex': latex,
                    })

        # If no structured parts found, extract plain text paragraphs
        if not body_parts:
            for p in el.find_all(['p', 'div'], recursive=False):
                if p.find('h1') or p.find('h2') or p.find('h3'):
                    continue
                text = _clean_text_with_math(p)
                if text.strip():
                    body_parts.append({'kind': 'text', 'content': text})

        slide_data['body_parts'] = body_parts

        # Extract speaker notes
        notes_el = el.find(class_=re.compile(r'note|speaker'))
        if notes_el:
            slide_data['notes'] = notes_el.get_text(strip=True)

        slides.append(slide_data)

    return slides


def _clean_text_with_math(element):
    """Extract text from an element, replacing math spans with LaTeX markers."""
    if element is None:
        return ''
    text = element.get_text(strip=True)
    # Clean up common HTML artifacts
    text = re.sub(r'\s+', ' ', text)
    return text


# ---------------------------------------------------------------------------
# PPTX Generation
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str):
    """Convert '#RRGGBB' to RGBColor."""
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 6:
        return RGBColor(
            int(hex_str[0:2], 16),
            int(hex_str[2:4], 16),
            int(hex_str[4:6], 16)
        )
    return RGBColor(0x33, 0x33, 0x33)


def _load_style_report(style_report_path):
    """Load and return style report JSON, or empty dict on failure."""
    if not style_report_path or not os.path.exists(style_report_path):
        return {}
    try:
        import json
        with open(style_report_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"  Warning: Could not load style report: {e}", file=sys.stderr)
        return {}


def build_pptx(slides_data, output_path, assets_dir=None,
               slide_width_in=DEFAULT_WIDTH, slide_height_in=DEFAULT_HEIGHT,
               style_report_path=None):
    """Build a PPTX presentation from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    # Load style report for logo positioning metadata
    style_report = _load_style_report(style_report_path)

    # Use blank layout for full control
    blank_layout = prs.slide_layouts[6]  # Blank

    # Check for template assets
    logo_path = None
    title_bg_path = None
    end_bg_path = None
    content_bg_path = None

    # Logo positioning from style report
    logo_left_in = None
    logo_top_in = None
    logo_width_in_val = 0.7

    if assets_dir:
        assets = Path(assets_dir)
        for candidate in ['logo.png', 'logo.jpg', 'logo.svg']:
            if (assets / candidate).exists():
                logo_path = str(assets / candidate)
                break
        for candidate in ['title_bg.png', 'title_bg.jpg']:
            if (assets / candidate).exists():
                title_bg_path = str(assets / candidate)
                break
        for candidate in ['end_bg.png', 'end_bg.jpg']:
            if (assets / candidate).exists():
                end_bg_path = str(assets / candidate)
                break
        for candidate in ['content_bg.png', 'content_bg.jpg']:
            if (assets / candidate).exists():
                content_bg_path = str(assets / candidate)
                break

    # Extract logo position from style report if available
    logo_candidates = style_report.get('logo_candidates', [])
    if logo_candidates and logo_path:
        top_candidate = logo_candidates[0]
        logo_left_in = round(top_candidate.get('left_emu', 0) / 914400, 2)
        logo_top_in = round(top_candidate.get('top_emu', 0) / 914400, 2)
        logo_width_in_val = top_candidate.get('width_in', 0.7)

    content_area_w = slide_width_in - MARGIN_LEFT - MARGIN_RIGHT
    content_area_h = slide_height_in - MARGIN_TOP - MARGIN_BOTTOM

    math_cache = {}  # Cache rendered math images

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        print(f"  Slide {idx + 1}/{len(slides_data)}: "
              f"{slide_data['heading'][:40]}...")

        # --- Background ---
        bg_path = slide_data.get('background_image')
        if not bg_path:
            if slide_data['type'] == 'title' and title_bg_path:
                bg_path = title_bg_path
            elif slide_data['type'] == 'end' and end_bg_path:
                bg_path = end_bg_path
            elif content_bg_path:
                bg_path = content_bg_path

        if bg_path and os.path.exists(bg_path):
            add_background_image(
                slide, bg_path,
                prs.slide_width, prs.slide_height
            )

        # --- Logo ---
        if logo_path:
            add_logo(slide, logo_path, 'top-right',
                     logo_width_in=logo_width_in_val,
                     slide_width_in=slide_width_in,
                     slide_height_in=slide_height_in,
                     exact_left_in=logo_left_in,
                     exact_top_in=logo_top_in)

        # --- Title slide ---
        if slide_data['type'] == 'title':
            _build_title_slide(slide, slide_data,
                               slide_width_in, slide_height_in)
            continue

        # --- End slide ---
        if slide_data['type'] == 'end':
            _build_end_slide(slide, slide_data,
                             slide_width_in, slide_height_in)
            continue

        # --- Content slide ---
        # Heading
        y_cursor = MARGIN_TOP
        if slide_data['heading']:
            txBox = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT), Inches(y_cursor - 0.1),
                Inches(content_area_w), Inches(0.6)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data['heading']
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            y_cursor += 0.7

        # Body parts
        for part in slide_data['body_parts']:
            remaining_h = slide_height_in - y_cursor - MARGIN_BOTTOM

            if remaining_h < 0.3:
                break  # No more room

            if part['kind'] == 'text':
                h = _add_text_part(slide, part['content'],
                                   MARGIN_LEFT, y_cursor,
                                   content_area_w, remaining_h)
                y_cursor += h + 0.1

            elif part['kind'] == 'bullet':
                h = _add_bullet_part(slide, part['items'],
                                     MARGIN_LEFT, y_cursor,
                                     content_area_w, remaining_h)
                y_cursor += h + 0.1

            elif part['kind'] in ('math_inline', 'math_display'):
                latex = part['latex']
                is_display = part['kind'] == 'math_display'

                # Check cache
                cache_key = (latex, is_display)
                if cache_key in math_cache:
                    img_buf, w_in, h_in = math_cache[cache_key]
                    img_buf.seek(0)
                else:
                    result = render_latex_to_png(
                        latex, display_mode=is_display
                    )
                    if result is None:
                        # Fallback: render as plain text
                        h = _add_text_part(
                            slide, latex,
                            MARGIN_LEFT, y_cursor,
                            content_area_w, remaining_h
                        )
                        y_cursor += h + 0.1
                        continue
                    img_buf, w_in, h_in = result
                    math_cache[cache_key] = (img_buf, w_in, h_in)
                    img_buf.seek(0)

                # Fit math image within content area
                max_math_w = content_area_w * 0.9
                max_math_h = min(remaining_h, 2.0)
                if w_in > max_math_w:
                    scale = max_math_w / w_in
                    w_in *= scale
                    h_in *= scale
                if h_in > max_math_h:
                    scale = max_math_h / h_in
                    w_in *= scale
                    h_in *= scale

                # Center math horizontally
                x = MARGIN_LEFT + (content_area_w - w_in) / 2
                try:
                    slide.shapes.add_picture(
                        img_buf,
                        Inches(x), Inches(y_cursor),
                        Inches(w_in), Inches(h_in)
                    )
                except Exception as e:
                    print(f"  Warning: Could not embed math image: {e}",
                          file=sys.stderr)

                y_cursor += h_in + 0.15

            elif part['kind'] == 'image':
                img_path = part['src']
                if not os.path.exists(img_path):
                    continue

                # Fit image within remaining content area
                max_img_w = content_area_w * 0.85
                max_img_h = min(remaining_h - 0.3, content_area_h * 0.55)
                if max_img_h < 0.5:
                    continue  # Not enough room

                w_in, h_in = fit_image_to_area(
                    img_path, max_img_w, max_img_h
                )

                # Center image horizontally
                x = MARGIN_LEFT + (content_area_w - w_in) / 2
                try:
                    slide.shapes.add_picture(
                        img_path,
                        Inches(x), Inches(y_cursor),
                        Inches(w_in), Inches(h_in)
                    )
                except Exception as e:
                    print(f"  Warning: Could not add image {img_path}: {e}",
                          file=sys.stderr)

                y_cursor += h_in + 0.15

                # Add caption if available
                if part.get('alt'):
                    h = _add_text_part(
                        slide, part['alt'],
                        MARGIN_LEFT, y_cursor,
                        content_area_w, 0.3,
                        font_size=10, italic=True,
                        alignment=PP_ALIGN.CENTER
                    )
                    y_cursor += h + 0.1

        # Speaker notes
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']

    prs.save(output_path)
    print(f"\nSaved PPTX to: {output_path}")
    print(f"  {len(slides_data)} slides, "
          f"{len(math_cache)} math expressions rendered")


def _build_title_slide(slide, data, slide_w, slide_h):
    """Build a title slide."""
    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(slide_h * 0.3),
        Inches(slide_w - 3), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['heading']
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if data.get('subheading'):
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), Inches(slide_h * 0.55),
            Inches(slide_w - 3), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = data['subheading']
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER


def _build_end_slide(slide, data, slide_w, slide_h):
    """Build an end/thank-you slide."""
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(slide_h * 0.35),
        Inches(slide_w - 3), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['heading'] or 'Thank You'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    p.alignment = PP_ALIGN.CENTER

    # Body text (e.g., contact info)
    for part in data['body_parts']:
        if part['kind'] == 'text':
            txBox2 = slide.shapes.add_textbox(
                Inches(1.5), Inches(slide_h * 0.55),
                Inches(slide_w - 3), Inches(1.0)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = part['content']
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.alignment = PP_ALIGN.CENTER
            break


def _add_text_part(slide, text, x, y, w, max_h,
                   font_size=18, italic=False, alignment=PP_ALIGN.LEFT):
    """Add a text paragraph to the slide. Returns the height used."""
    h = min(max_h, 0.5)
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.italic = italic
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.alignment = alignment
    return h


def _add_bullet_part(slide, items, x, y, w, max_h):
    """Add a bulleted list to the slide. Returns the height used."""
    line_h = 0.35
    h = min(max_h, line_h * len(items) + 0.1)

    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.level = 0
        # Add bullet character
        p.text = f"  \u2022  {item}"

    return h


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description='Convert HTML slides to PPTX with math & figure support'
    )
    parser.add_argument('html_file', help='Path to the HTML presentation')
    parser.add_argument('--output', '-o', default=None,
                        help='Output PPTX path (default: same name as input)')
    parser.add_argument('--assets', default=None,
                        help='Path to assets directory (backgrounds, logos)')
    parser.add_argument('--width', type=float, default=DEFAULT_WIDTH,
                        help=f'Slide width in inches (default: {DEFAULT_WIDTH})')
    parser.add_argument('--height', type=float, default=DEFAULT_HEIGHT,
                        help=f'Slide height in inches (default: {DEFAULT_HEIGHT})')
    parser.add_argument('--style-report', default=None,
                        help='Path to style_report.json for logo positioning and metadata')
    args = parser.parse_args()

    html_path = Path(args.html_file)
    if not html_path.exists():
        print(f"Error: {html_path} not found", file=sys.stderr)
        sys.exit(1)

    output_path = args.output or str(html_path.with_suffix('.pptx'))
    assets_dir = args.assets or str(html_path.parent / 'assets')

    print(f"Converting: {html_path}")
    print(f"Output:     {output_path}")
    print(f"Assets:     {assets_dir}")
    print(f"Slide size: {args.width}\" x {args.height}\"")
    print()

    # Parse HTML
    print("Parsing HTML slides...")
    slides = parse_html_slides(str(html_path))

    if not slides:
        print("Error: No slides found in HTML file.", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(slides)} slides.\n")

    # Build PPTX
    print("Generating PPTX...")
    build_pptx(
        slides, output_path, assets_dir,
        slide_width_in=args.width, slide_height_in=args.height,
        style_report_path=args.style_report
    )


if __name__ == '__main__':
    main()
