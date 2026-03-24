#!/usr/bin/env python3
"""
extract_pptx_template.py — Reusable PPTX template analysis & asset extraction.

Extracts backgrounds, logos, colors, fonts, and layout metadata from a .pptx
file, producing a JSON style report and saving all visual assets to an output
directory. Critically, it detects and applies OOXML image transforms (grayscale,
duotone, etc.) so that saved assets match the actual PowerPoint appearance.

Usage:
    python extract_pptx_template.py <pptx_path> --output-dir <assets_dir>

Output:
    <assets_dir>/
        title_bg.png          — title-slide background (with transforms applied)
        content_header.png    — content-slide header bar (with transforms applied)
        end_bg.png            — end/thank-you-slide background (with transforms applied)
        logo_*.png            — any extracted logos
        style_report.json     — full style metadata

Persisted at: ~/.claude/skills/paper-to-slides/extract_pptx_template.py
"""

import argparse
import json
import os
import sys
import zipfile
from pathlib import Path

# --- Dependency management ---------------------------------------------------

def ensure_import(package_name, pip_name=None):
    """Import a package, auto-installing via pip if missing."""
    try:
        return __import__(package_name)
    except ImportError:
        import subprocess
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", pip_name or package_name],
            stdout=subprocess.DEVNULL,
        )
        return __import__(package_name)


# --- OOXML image-transform detection ----------------------------------------

# PowerPoint can apply visual transforms to embedded images via child elements
# inside <a:blip>. The raw image blob from python-pptx does NOT include these;
# we must detect and replicate them in Pillow.
#
# Known transforms (OOXML DrawingML):
#   <a:grayscl/>         — convert to grayscale
#   <a:duotone>          — map to two-color gradient
#   <a:biLevel thresh/>  — black & white with threshold
#   <a:lum bright= contrast=/>  — brightness / contrast adjustment
#   <a:alphaModFix amt=/>— change opacity
#   <a:clrRepl>          — replace a color
#   <a:clrChange>        — remap one color to another
#
# The most common one encountered in real templates is <a:grayscl/>.

def _ns(tag):
    """Expand a short namespace prefix into full Clark notation."""
    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    prefix, local = tag.split(":")
    return f"{{{nsmap[prefix]}}}{local}"


def detect_blip_transforms(blip_element):
    """
    Inspect an <a:blip> element for image transforms.

    Returns a list of transform dicts, e.g.:
        [{"type": "grayscale"}]
        [{"type": "duotone", "colors": [...]}]
        [{"type": "lum", "bright": 20000, "contrast": 0}]
    """
    transforms = []
    if blip_element is None:
        return transforms

    for child in blip_element:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "grayscl":
            transforms.append({"type": "grayscale"})

        elif tag == "biLevel":
            thresh = child.get("thresh", "50000")
            transforms.append({"type": "bilevel", "threshold": int(thresh)})

        elif tag == "lum":
            transforms.append({
                "type": "luminance",
                "bright": int(child.get("bright", "0")),
                "contrast": int(child.get("contrast", "0")),
            })

        elif tag == "alphaModFix":
            transforms.append({
                "type": "alpha",
                "amount": int(child.get("amt", "100000")),
            })

        elif tag == "duotone":
            colors = []
            for clr_child in child:
                clr_tag = clr_child.tag.split("}")[-1]
                if clr_tag == "srgbClr":
                    colors.append(clr_child.get("val"))
                elif clr_tag == "schemeClr":
                    colors.append(f"scheme:{clr_child.get('val')}")
            transforms.append({"type": "duotone", "colors": colors})

    return transforms


def apply_pil_transforms(img, transforms):
    """
    Apply detected OOXML transforms to a PIL Image.

    Args:
        img: PIL.Image instance (RGB or RGBA)
        transforms: list of transform dicts from detect_blip_transforms()

    Returns:
        New PIL.Image with transforms applied (never mutates the original).
    """
    from PIL import ImageOps, ImageEnhance

    result = img.copy()

    for t in transforms:
        if t["type"] == "grayscale":
            # Convert to grayscale, then back to RGB so it can be saved as PNG
            result = ImageOps.grayscale(result).convert("RGB")

        elif t["type"] == "bilevel":
            # Threshold: values in 1/100000 units, 50000 = 50%
            thresh_pct = t.get("threshold", 50000) / 1000.0  # → 0–100
            thresh_val = int(thresh_pct * 255 / 100)
            gray = ImageOps.grayscale(result)
            result = gray.point(lambda p: 255 if p >= thresh_val else 0).convert("RGB")

        elif t["type"] == "luminance":
            # bright/contrast in 1/1000 units (e.g., 20000 = +20%)
            bright_factor = 1.0 + t.get("bright", 0) / 100000.0
            contrast_factor = 1.0 + t.get("contrast", 0) / 100000.0
            result = ImageEnhance.Brightness(result).enhance(bright_factor)
            result = ImageEnhance.Contrast(result).enhance(contrast_factor)

        elif t["type"] == "alpha":
            # Amount in 1/100000 units (100000 = fully opaque)
            alpha_pct = t.get("amount", 100000) / 100000.0
            if result.mode != "RGBA":
                result = result.convert("RGBA")
            r, g, b, a = result.split()
            from PIL import ImageChops
            a = a.point(lambda p: int(p * alpha_pct))
            result = Image.merge("RGBA", (r, g, b, a))      # type: ignore

        # duotone is complex and rarely used; skip with a warning
        elif t["type"] == "duotone":
            print(f"  [WARN] Duotone transform detected but not applied (colors: {t['colors']}). "
                  "Manual adjustment may be needed.", file=sys.stderr)

    return result


# --- Aspect ratio classification --------------------------------------------

def classify_aspect_ratio(width_emu, height_emu):
    """Classify slide dimensions as '16:9' or '4:3' (with tolerance).

    Common EMU values:
        16:9  →  12192000 x 6858000  (13.333" x 7.5")
        4:3   →  9144000  x 6858000  (10" x 7.5")

    Returns:
        A dict with 'label' ('16:9' or '4:3'), 'ratio' (float),
        and 'other_label' (the opposite ratio label).
    """
    ratio = width_emu / height_emu if height_emu else 0

    # 16:9 = 1.7778, 4:3 = 1.3333
    KNOWN_RATIOS = [
        ("16:9", 16 / 9),
        ("4:3", 4 / 3),
    ]

    best_label = None
    best_distance = float("inf")
    for label, target in KNOWN_RATIOS:
        dist = abs(ratio - target)
        if dist < best_distance:
            best_distance = dist
            best_label = label

    # Only accept if within 5% tolerance
    if best_distance > 0.1:
        best_label = f"custom ({ratio:.4f})"

    other_label = None
    if best_label == "16:9":
        other_label = "4:3"
    elif best_label == "4:3":
        other_label = "16:9"

    return {
        "label": best_label,
        "ratio": round(ratio, 4),
        "other_label": other_label,
    }


# --- Logo detection ----------------------------------------------------------

def identify_logo_candidates(prs, output_dir):
    """Scan slide master and layouts for picture shapes that are likely logos.

    Scores each candidate by heuristics:
      - Size < 15% of slide area (+3)
      - Position in corner region (+2)
      - Shape name contains "logo"/"emblem"/"badge" (+2)
      - Present on master (inherited by all slides) (+1)
      - Aspect ratio near 1:1 or known logo proportions (+1)

    Returns a sorted list of candidate dicts (highest score first).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io

    slide_area = prs.slide_width * prs.slide_height
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    candidates = []

    def _score_shape(shape, source_label, is_master):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None

        score = 0
        shape_area = shape.width * shape.height

        # Size < 15% of slide area
        if shape_area < slide_area * 0.15:
            score += 3

        # Corner region (top/bottom 25%, left/right 25%)
        cx = shape.left + shape.width / 2
        cy = shape.top + shape.height / 2
        in_left_right = (cx < slide_w * 0.25) or (cx > slide_w * 0.75)
        in_top_bottom = (cy < slide_h * 0.25) or (cy > slide_h * 0.75)
        if in_left_right and in_top_bottom:
            score += 2

        # Name contains logo-related keywords
        name_lower = shape.name.lower()
        if any(kw in name_lower for kw in ("logo", "emblem", "badge", "icon")):
            score += 2

        # Present on master
        if is_master:
            score += 1

        # Aspect ratio near 1:1 (square logos)
        if shape.width > 0 and shape.height > 0:
            ar = shape.width / shape.height
            if 0.5 <= ar <= 2.0:
                score += 1

        # Skip very large images (likely backgrounds, not logos)
        if shape_area > slide_area * 0.5:
            return None

        # Save the image
        ext = shape.image.content_type.split("/")[-1]
        if ext == "jpeg":
            ext = "jpg"
        safe_name = shape.name.replace(" ", "_").replace("/", "_")
        fname = f"logo_candidate_{source_label}_{safe_name}.{ext}"
        fpath = os.path.join(output_dir, fname)

        img = Image.open(io.BytesIO(shape.image.blob))

        # Apply transforms if any
        pic_el = shape._element
        blip_el = pic_el.find(".//" + _ns("a:blip"))
        transforms = detect_blip_transforms(blip_el)
        if transforms:
            img = apply_pil_transforms(img, transforms)
            fname = fname.rsplit(".", 1)[0] + ".png"
            fpath = os.path.join(output_dir, fname)

        img.save(fpath)

        # Determine position label
        pos_label = ""
        if cy < slide_h * 0.5:
            pos_label += "top-"
        else:
            pos_label += "bottom-"
        if cx < slide_w * 0.5:
            pos_label += "left"
        else:
            pos_label += "right"

        return {
            "file": fname,
            "name": shape.name,
            "score": score,
            "source": source_label,
            "is_master": is_master,
            "position": pos_label,
            "left_emu": shape.left,
            "top_emu": shape.top,
            "width_emu": shape.width,
            "height_emu": shape.height,
            "width_in": round(shape.width / 914400, 2),
            "height_in": round(shape.height / 914400, 2),
        }

    # Scan masters
    for master in prs.slide_masters:
        for shape in master.shapes:
            c = _score_shape(shape, "master", is_master=True)
            if c:
                candidates.append(c)

    # Scan layouts
    for i, layout in enumerate(prs.slide_layouts):
        for shape in layout.shapes:
            c = _score_shape(shape, f"layout{i}", is_master=False)
            if c:
                candidates.append(c)

    # Sort by score descending
    candidates.sort(key=lambda c: c["score"], reverse=True)
    return candidates


def generate_logo_preview(candidates, output_dir):
    """Create an HTML preview file showing all logo candidates in a grid.

    Returns the path to the generated HTML file.
    """
    if not candidates:
        return None

    html_parts = [
        "<!DOCTYPE html>",
        "<html><head><meta charset='utf-8'>",
        "<title>Logo Candidates</title>",
        "<style>",
        "body { font-family: system-ui, sans-serif; max-width: 1000px; margin: 2rem auto; padding: 0 1rem; }",
        ".grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(250px, 1fr)); gap: 1.5rem; }",
        ".card { border: 1px solid #ddd; border-radius: 8px; padding: 1rem; text-align: center; }",
        ".card img { max-width: 150px; max-height: 150px; object-fit: contain; margin: 0.5rem 0; }",
        ".score { font-size: 1.5rem; font-weight: bold; color: #2563eb; }",
        ".meta { font-size: 0.85rem; color: #666; }",
        "h1 { color: #111; }",
        "</style></head><body>",
        "<h1>Logo Candidates</h1>",
        f"<p>Found {len(candidates)} candidate(s). Higher score = more likely a logo.</p>",
        "<div class='grid'>",
    ]

    for i, c in enumerate(candidates):
        html_parts.append(f"""
        <div class='card'>
            <div class='score'>Score: {c['score']}</div>
            <img src='{c['file']}' alt='{c['name']}'>
            <div><strong>#{i + 1}: {c['name']}</strong></div>
            <div class='meta'>
                Source: {c['source']}<br>
                Position: {c['position']}<br>
                Size: {c['width_in']}" x {c['height_in']}"
            </div>
        </div>""")

    html_parts.append("</div></body></html>")

    preview_path = os.path.join(output_dir, "_logo_preview.html")
    with open(preview_path, "w", encoding="utf-8") as f:
        f.write("\n".join(html_parts))

    print(f"  Logo preview saved to: {preview_path}")
    return preview_path


# --- Main extraction ---------------------------------------------------------

def extract_template(pptx_path, output_dir):
    """
    Extract all visual assets and style metadata from a PPTX template.

    Returns a dict with the full style report.
    """
    ensure_import("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io
    import lxml.etree as ET     # type: ignore[import]  # lxml does not have type hints

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    ar_info = classify_aspect_ratio(prs.slide_width, prs.slide_height)

    report = {
        "source": str(pptx_path),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slide_width_in": round(prs.slide_width / 914400, 2),
        "slide_height_in": round(prs.slide_height / 914400, 2),
        "aspect_ratio": round(prs.slide_width / prs.slide_height, 4),
        "aspect_ratio_label": ar_info["label"],
        "aspect_ratio_other": ar_info["other_label"],
        "num_slides": len(prs.slides),
        "num_layouts": len(prs.slide_layouts),
        "theme_colors": {},
        "fonts": set(),
        "layouts": [],
        "master_images": [],
        "slide_texts": [],
        "image_transforms_detected": [],
    }

    # --- Extract theme colors ------------------------------------------------
    try:
        with zipfile.ZipFile(pptx_path) as z:
            with z.open("ppt/theme/theme1.xml") as f:
                tree = ET.parse(f)
                root = tree.getroot()
                clr_scheme = root.find(".//" + _ns("a:clrScheme"))
                if clr_scheme is not None:
                    report["theme_name"] = clr_scheme.get("name", "unknown")
                    for child in clr_scheme:
                        color_name = child.tag.split("}")[-1]
                        for color_el in child:
                            color_tag = color_el.tag.split("}")[-1]
                            val = color_el.get("val", color_el.get("lastClr", ""))
                            report["theme_colors"][color_name] = {
                                "type": color_tag,
                                "value": val,
                            }
    except Exception as e:
        print(f"  [WARN] Could not extract theme: {e}", file=sys.stderr)

    # --- Helper: extract images from a shape collection ----------------------
    def _extract_images(shapes, prefix, record_list):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = shape.image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                safe_name = shape.name.replace(" ", "_").replace("/", "_")
                fname = f"{prefix}_{safe_name}.{ext}"
                fpath = os.path.join(output_dir, fname)

                # Detect transforms on the blip element
                pic_el = shape._element
                blip_el = pic_el.find(".//" + _ns("a:blip"))
                transforms = detect_blip_transforms(blip_el)

                # Save raw image
                img = Image.open(io.BytesIO(shape.image.blob))

                # Apply any detected transforms
                if transforms:
                    report["image_transforms_detected"].append({
                        "image": fname,
                        "transforms": transforms,
                    })
                    img = apply_pil_transforms(img, transforms)
                    # Change extension to PNG since we've processed it
                    fname = fname.rsplit(".", 1)[0] + ".png"
                    fpath = os.path.join(output_dir, fname)

                img.save(fpath)
                record_list.append({
                    "file": fname,
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "transforms": transforms,
                })
                print(f"  Extracted: {fname} ({shape.width}x{shape.height} EMU)"
                      + (f" [transforms: {[t['type'] for t in transforms]}]" if transforms else ""))

            # Collect fonts from text shapes
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            report["fonts"].add(run.font.name)

    # --- Extract from slide master -------------------------------------------
    print("Extracting from slide master...")
    master_images = []
    for master in prs.slide_masters:
        _extract_images(master.shapes, "master", master_images)

        # Check master for decorative shapes (bars, rectangles)
        for shape in master.shapes:
            if not hasattr(shape, "fill"):
                continue
            sp_el = shape._element
            solid_fill = sp_el.find(".//" + _ns("a:solidFill"))
            if solid_fill is not None:
                scheme_clr = solid_fill.find(_ns("a:schemeClr"))
                srgb_clr = solid_fill.find(_ns("a:srgbClr"))
                if scheme_clr is not None:
                    color_info = {"scheme": scheme_clr.get("val")}
                    lum_mod = scheme_clr.find(_ns("a:lumMod"))
                    if lum_mod is not None:
                        color_info["lumMod"] = int(lum_mod.get("val"))
                elif srgb_clr is not None:
                    color_info = {"srgb": srgb_clr.get("val")}
                else:
                    color_info = {}

                if color_info and shape.width > prs.slide_width * 0.5:
                    master_images.append({
                        "type": "decorative_bar",
                        "name": shape.name,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "fill": color_info,
                    })

    report["master_images"] = master_images

    # --- Extract from each layout --------------------------------------------
    print("Extracting from slide layouts...")
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {"index": i, "name": layout.name, "images": []}
        _extract_images(layout.shapes, f"layout{i}", layout_info["images"])
        report["layouts"].append(layout_info)

    # --- Extract text content from actual slides -----------------------------
    print("Extracting slide text...")
    for idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        font_info = {}
                        for run in para.runs:
                            f = run.font
                            try:
                                color_val = str(f.color.rgb) if f.color and f.color.type else None
                            except (AttributeError, TypeError):
                                color_val = None
                            font_info = {
                                "font": f.name,
                                "size": f.size,
                                "bold": f.bold,
                                "color": color_val,
                            }
                            if f.name:
                                report["fonts"].add(f.name)
                            break
                        slide_texts.append({"text": text[:200], "font": font_info})
        report["slide_texts"].append({
            "slide": idx + 1,
            "texts": slide_texts,
        })

    # --- Assign canonical names to key assets --------------------------------
    # Try to identify title_bg, content_header, end_bg from layout/master images
    canonical = {}
    if report["layouts"] and report["layouts"][0]["images"]:
        first_layout_img = report["layouts"][0]["images"][0]
        src = os.path.join(output_dir, first_layout_img["file"])
        dst = os.path.join(output_dir, "title_bg.png")
        if os.path.exists(src):
            import shutil
            shutil.copy2(src, dst)
            canonical["title_bg"] = "title_bg.png"

    for img_info in master_images:
        if isinstance(img_info, dict) and "file" in img_info:
            # Narrow/wide image at top → likely header bar
            if img_info.get("height", 0) < prs.slide_height * 0.15:
                src = os.path.join(output_dir, img_info["file"])
                dst = os.path.join(output_dir, "content_header.png")
                if os.path.exists(src):
                    import shutil
                    shutil.copy2(src, dst)
                    canonical["content_header"] = "content_header.png"

    report["canonical_assets"] = canonical

    # --- Detect logo candidates -----------------------------------------------
    print("Scanning for logo candidates...")
    logo_candidates = identify_logo_candidates(prs, output_dir)
    report["logo_candidates"] = logo_candidates
    if logo_candidates:
        preview_path = generate_logo_preview(logo_candidates, output_dir)
        report["logo_preview"] = os.path.basename(preview_path) if preview_path else None
        print(f"  Found {len(logo_candidates)} logo candidate(s)")
    else:
        report["logo_preview"] = None
        print("  No logo candidates found")

    # Convert set to list for JSON serialization
    report["fonts"] = sorted(report["fonts"])

    # --- Write report --------------------------------------------------------
    report_path = os.path.join(output_dir, "style_report.json")
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False, default=str)
    print(f"\nStyle report saved to: {report_path}")

    return report


# --- CLI entry point ---------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Extract visual assets and style metadata from a PPTX template."
    )
    parser.add_argument("pptx_path", help="Path to the .pptx template file")
    parser.add_argument(
        "--output-dir", "-o",
        default="assets",
        help="Directory to save extracted assets (default: assets/)",
    )
    args = parser.parse_args()

    if not os.path.isfile(args.pptx_path):
        print(f"Error: file not found: {args.pptx_path}", file=sys.stderr)
        sys.exit(1)

    report = extract_template(args.pptx_path, args.output_dir)

    # Print summary
    print(f"\n{'='*60}")
    print(f"Template: {report['source']}")
    print(f"Dimensions: {report['slide_width_in']}\" x {report['slide_height_in']}\" "
          f"(aspect {report['aspect_ratio']}, {report['aspect_ratio_label']})")
    print(f"Theme: {report.get('theme_name', 'unknown')}")
    print(f"Fonts: {', '.join(report['fonts']) or 'none detected'}")
    if report["image_transforms_detected"]:
        print(f"\n⚠ Image transforms detected and applied:")
        for t in report["image_transforms_detected"]:
            print(f"  {t['image']}: {[x['type'] for x in t['transforms']]}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
